import boto3
import sys
import os
import json
import io
import logging
import re
from strands import Agent
from pptx import Presentation
from pptx.util import Inches
import pdfkit
import pypandoc
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
import time
from PIL import Image

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

def clean_html_output(raw_output, is_presentation=False):
    try:
        if hasattr(raw_output, 'output'):
            output = str(raw_output.output)
        elif hasattr(raw_output, 'text'):
            output = str(raw_output.text)
        else:
            output = str(raw_output)

        with open("raw_agent_output.txt", "w", encoding="utf-8") as f:
            f.write(output)
        logger.info("Saved raw Agent output to raw_agent_output.txt")

        output = re.sub(r'```html\s*\n?', '', output)
        output = re.sub(r'```\s*$', '', output, flags=re.MULTILINE)
        output = re.sub(r'```[a-zA-Z]*\s*\n?', '', output)
        output = output.replace('```', '')

        output = output.strip()
        if not output.startswith(('<!DOCTYPE html', '<html')):
            logger.error("Output does not start with valid HTML.")
            raise ValueError("Output is not valid HTML")

        if '<html' not in output.lower() or '</html>' not in output.lower():
            logger.error("Invalid HTML structure detected.")
            raise ValueError("Output is not valid HTML")

        return output
    except Exception as e:
        logger.error(f"Failed to clean HTML output: {str(e)}")
        raise

def process_screenshot(image_path, target_width=1280, target_height=720):
    try:
        with Image.open(image_path) as img:
            img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)
            processed_path = image_path.replace(".png", "_processed.png")
            img.save(processed_path, quality=95)
            return processed_path
    except Exception as e:
        logger.error(f"Failed to process screenshot {image_path}: {str(e)}")
        return image_path

# ---------- STEP 1: READ PAYLOAD ----------
try:
    with open("input.json", "r", encoding="utf-8") as f:
        use_case = json.load(f)
    use_case_name = use_case.get("use_case_name", "Untitled Use Case")
    description = use_case.get("description", "")
    project_id = use_case.get("project_id", "report")
    user_id = use_case.get("user_id", "default_user")
    s3_bucket = "qubitz-customer-prod-parv"
    s3_prefix = f"{user_id}/{project_id}/"
except Exception as e:
    logger.error(f"Failed to read input.json: {str(e)}")
    raise

agent = Agent()

# ---------- STEP 2: GENERATE PRESENTATION FIRST ----------
ppt_prompt = f"""
Generate a clean, modern, fully styled HTML5 presentation (not markdown) with 8-10 slides.

🎯 Objectives:
- Designed for executives
- No emojis or unnecessary decorative icons
- Each slide should be a <div class="slide"> and same size

📘 Theme:
- Business use case: "{use_case_name}"
- Description: {description}

📐 Structure:
1. Title & Overview
2. Current Challenges
3. Key Features
4. Workflow
5. Data Enrichment
6. AI Personalization
7. Implementation Timeline
8. Metrics / ROI
9. Business Impact
10. Call to Action

🎨 Visual & Styling (embed with inline CSS):
- Use clean enterprise fonts (Segoe UI or similar)
- Two-column layouts, metric cards, call-to-action banners, timelines, hover/transition effects
- Consistent padding and spacing
- Avoid JS, images, markdown, or emojis
- Include professional inline CSS for each component (cards, bullets, boxes, etc.)

📝 Notes:
- Each section must be self-contained inside a <div class="slide">
- Use consistent slide dimensions
- Keep animations subtle (fade-in, transitions if CSS-only)

📏 Constraints:
- Each <div class="slide"> must be exactly 720px in height
- Do NOT use scrollbars or overflow
- Keep content within bounds using layout techniques: columns, font size, spacing

"""

try:
    ppt_response = agent(ppt_prompt)
    ppt_html = clean_html_output(ppt_response, is_presentation=True)
    ppt_filename = f"{project_id}_presentation.html"
    with open(ppt_filename, "w", encoding="utf-8") as f:
        f.write(ppt_html)
    logger.info(f"HTML presentation generated: {ppt_filename}")
except Exception as e:
    logger.error(f"Presentation HTML generation failed: {str(e)}")
    raise

# ---------- STEP 3: CONVERT TO PPTX ----------
def html_to_image_pptx(html_path):
    try:
        base = os.path.splitext(html_path)[0]
        output_pptx = base + ".pptx"
        screenshot_folder = "slides"
        os.makedirs(screenshot_folder, exist_ok=True)

        chrome_options = Options()
        chrome_options.add_argument("--headless=new")
        chrome_options.add_argument("--window-size=1280,720")
        chrome_options.add_argument("--disable-gpu")
        chrome_options.add_argument("--no-sandbox")
        chrome_options.add_argument("--force-device-scale-factor=1")
        chrome_options.add_argument("--high-dpi-support=1")

        driver = webdriver.Chrome(options=chrome_options)
        file_url = f"file:///{os.path.abspath(html_path).replace('\\', '/')}"
        driver.get(file_url)
        time.sleep(3)

        driver.execute_cdp_cmd('Emulation.setDeviceMetricsOverride', {
            "width": 1280,
            "height": 720,
            "deviceScaleFactor": 1,
            "mobile": False
        })

        slides = driver.find_elements(By.CLASS_NAME, "slide")
        image_paths = []
        for i, _ in enumerate(slides):
            driver.execute_script("""
                document.querySelectorAll('.slide').forEach((s, i2) => s.style.display = (i2 === arguments[0] ? 'block' : 'none'));
            """, i)
            time.sleep(1.5)
            path = os.path.join(screenshot_folder, f"slide_{i+1}.png")
            driver.get_screenshot_as_file(path)
            processed_path = process_screenshot(path)
            image_paths.append(processed_path)

        driver.quit()

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        blank_layout = prs.slide_layouts[6]
        for img_path in image_paths:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        prs.save(output_pptx)
        logger.info(f"PPTX created: {output_pptx}")
    except Exception as e:
        logger.error(f"PPTX generation failed: {str(e)}")
        raise

html_to_image_pptx(ppt_filename)

# ---------- STEP 4: GENERATE WORD REPORT ----------
word_prompt = f"""
Generate a formal HTML business report for Word and PDF, styled like a professional report. Content should focus on the use case: {use_case_name}.
- Output ONLY pure HTML with semantic tags: <h1>-<h3>, <p>, <ul>/<ol>, avoid <table>
- Do NOT include Markdown code fences
- Title: Business Strategy for {use_case_name}
- Structure:
  - Introduction: Summarize opportunity and potential value
  - Market Context: Industry challenges and digital trends
  - AI Applications: 5-8 use cases with ROI metrics
  - Implementation Roadmap: 3 phases (1-6, 6-12, 12-24 months)
  - Conclusion: Recap value and urgency
- Inline CSS: Arial, 14px, A4 width (595px), proper margins
❌ Do NOT include any headers, footers, copyright, page numbers, or boilerplate like “Confidential”, “Page 1 of 1”, or © statements.
"""
def upload_to_s3(local_path, bucket, key):
    try:
        s3 = boto3.client('s3')
        s3.upload_file(
            local_path, bucket, key,
            ExtraArgs={'ACL': 'public-read'}  # 👈 This makes the file publicly accessible
        )
        s3_url = f"https://{bucket}.s3.amazonaws.com/{key}"
        logger.info(f"Uploaded {local_path} to S3: {s3_url}")
        return s3_url
    except Exception as e:
        logger.error(f"Failed to upload {local_path} to S3: {str(e)}")
        return None
try:
    word_response = agent(word_prompt)
    word_html = clean_html_output(word_response, is_presentation=False)
    word_filename = f"{project_id}_report.html"
    with open(word_filename, "w", encoding="utf-8") as f:
        f.write(word_html)
    logger.info(f"HTML report generated: {word_filename}")

    pdf_path = word_filename.replace(".html", ".pdf")
    pdfkit.from_file(word_filename, pdf_path, options={
        'page-size': 'A4',
        'encoding': 'UTF-8',
        'enable-local-file-access': '',
        'margin-top': '0.75in',
        'margin-bottom': '0.75in',
        'margin-left': '0.75in',
        'margin-right': '0.75in'
    })
    logger.info(f"PDF created: {pdf_path}")

    docx_path = word_filename.replace(".html", ".docx")
    pypandoc.convert_file(word_filename, 'docx', outputfile=docx_path)
    logger.info(f"DOCX created: {docx_path}")
    
    # Upload output files to S3
    pptx_file = ppt_filename.replace(".html", ".pptx")
    pdf_file = word_filename.replace(".html", ".pdf")
    docx_file = word_filename.replace(".html", ".docx")

    pptx_url = upload_to_s3(pptx_file, s3_bucket, s3_prefix + os.path.basename(pptx_file))
    pdf_url  = upload_to_s3(pdf_file, s3_bucket, s3_prefix + os.path.basename(pdf_file))
    docx_url = upload_to_s3(docx_file, s3_bucket, s3_prefix + os.path.basename(docx_file))

    print("S3 URLs:")
    print(f" - PPTX: {pptx_url}")
    print(f" - PDF : {pdf_url}")
    print(f" - DOCX: {docx_url}")

except Exception as e:
    logger.error(f"Word/PDF/DOCX generation failed: {str(e)}")
    raise

print("✅ Done!")
